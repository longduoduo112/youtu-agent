import json
import logging
import re
from pathlib import Path
from typing import Any

import yaml
from ppt_template_model import PageConfig
from pptx import Presentation
from utils import delete_slide_range, duplicate_slide, move_slide


def fill_template_with_yaml_config(template_path, output_path, json_data, yaml_config: dict[str, Any]):
    page_config = PageConfig(yaml_config)
    prs = Presentation(template_path)
    data = json.loads(json_data)
    slides_data = data.get("slides", [])

    if not isinstance(slides_data, list):
        raise ValueError("JSON data must contain a 'slides' list")

    for slide_data in slides_data:
        slide_type = slide_data.get("type")
        if not slide_type:
            raise ValueError("Slide data must contain a 'type' field")

        template_index = page_config.type_map.get(slide_type)
        if template_index is None or template_index >= len(prs.slides):
            raise ValueError("No template found for slide type '%s'", slide_type)

        template_slide = prs.slides[template_index]
        if slide_type in ("title_page", "acknowledgement_page"):
            target_slide = template_slide
        else:
            target_slide = duplicate_slide(prs, template_slide)

        page_config.render(target_slide, slide_data)

    # get title page
    title_pages_idx = page_config.type_map.get("title_page")
    acknowledgement_pages_idx = page_config.type_map.get("acknowledgement_page")
    max_idx = max(list(page_config.type_map.values()))
    # move title page to the first
    move_slide(prs, title_pages_idx, 0)
    # move acknowledgement page to the last
    move_slide(prs, acknowledgement_pages_idx, len(prs.slides) - 1)
    # remove all
    delete_slide_range(prs, range(1, max_idx))
    prs.save(output_path)


def extract_json(content):
    """
    Extract the json data from the given content.
    """
    # extract content within "```json" and "```"
    pattern = r"```json(.*?)```"
    match = re.search(pattern, content, re.DOTALL)
    if match:
        return match.group(1)
    return content


if __name__ == "__main__":
    import argparse
    import datetime

    parser = argparse.ArgumentParser()
    parser.add_argument("-t", "--template", type=str, default="templates")
    parser.add_argument("-n", "--template_name", type=str, default="0")
    default_output_filename = f"output-{datetime.datetime.now().strftime('%Y-%m-%d-%H-%M-%S')}.pptx"
    parser.add_argument("-o", "--output", type=str, default=default_output_filename)
    parser.add_argument("-i", "--input", type=str, required=True)
    parser.add_argument("--cache_dir", type=str, default=".temp")
    parser.add_argument("--yaml_config", type=str, default="yaml_example.yaml")
    args = parser.parse_args()

    # set env var UTU_PPT_CACHE_DIR
    import os

    os.environ["UTU_PPT_CACHE_DIR"] = args.cache_dir

    logging.basicConfig(level=logging.INFO)
    template = args.template
    output = args.output
    input_json = args.input
    with open(input_json) as f:
        content = f.read()
    json_data = extract_json(content)
    if not json_data:
        raise ValueError("No JSON data found in input file")
    with open(args.yaml_config) as f:
        yaml_config = yaml.safe_load(f)
    template_yaml = Path(template) / args.template_name / f"{args.template_name}.yaml"
    if not template_yaml.exists():
        template_config = {}
    else:
        with open(template_yaml) as f:
            template_config = yaml.safe_load(f)
    # merge
    yaml_config.update(template_config)

    print(json.dumps(yaml_config, indent=4))

    template_path = Path(template) / args.template_name / f"{args.template_name}.pptx"

    fill_template_with_yaml_config(template_path, output, json_data, yaml_config)
