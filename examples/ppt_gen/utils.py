import copy
import logging
import random
from copy import deepcopy

import matplotlib
from pptx.enum.shapes import MSO_SHAPE_TYPE

# rgb colors from a color scheme
_color_palette = [
    (178, 34, 34),  # Brick Red
    (46, 139, 87),  # Sea Green
    (70, 130, 180),  # Steel Blue
    (210, 180, 140),  # Tan
    (147, 112, 219),  # Medium Purple
    (255, 165, 0),  # Orange (desaturated)
    (72, 209, 204),  # Medium Turquoise
    (205, 92, 92),  # Indian Red
    (106, 90, 205),  # Slate Blue
    (238, 130, 238),  # Violet
    (60, 179, 113),  # Medium Sea Green
    (100, 149, 237),  # Cornflower Blue
    (218, 165, 32),  # Goldenrod
    (199, 21, 133),  # Medium Violet Red
    (65, 105, 225),  # Royal Blue
]


# 常见 content-type -> reltype 映射（可扩展）
CONTENT_TYPE_TO_RELTYPE = {
    "image/jpeg": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
    "image/png": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
    "image/gif": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
    "application/vnd.openxmlformats-officedocument.drawingml.chart+xml": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart",
    "application/vnd.openxmlformats-officedocument.oleObject": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/oleObject",
    # 根据需要添加更多映射
}


def inspect_ppt(prs):
    """
    Inspect the given presentation.
    """
    for slide in prs.slides:
        print(f"==={slide.slide_layout.name}===")
        inspect_slide(slide)


def inspect_slide(slide):
    """
    Inspect the given slide layout.
    """

    def _inspect_shape_list(shapes, indent=0):
        for shape in shapes:
            shape_type_str = shape.shape_type.name
            print(" " * indent + shape.name + "(" + shape_type_str + ")")
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                _inspect_shape_list(shape.shapes, indent + 2)

    _inspect_shape_list(slide.shapes)


def to_svg(slide, prs, svg_filename="test.svg"):
    """
    Convert the given slide layout to an SVG file.
    """

    def _to_svg_box(shapes, svg_box):
        for shape in shapes:
            svg_box.append(
                {
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                    "shape_type": shape.shape_type,
                }
            )
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                _to_svg_box(shape.shapes, svg_box)

    def render_svg(svg_box):
        width = prs.slide_width.inches
        height = prs.slide_height.inches
        fig, ax = matplotlib.pyplot.subplots(figsize=(width, height))
        ax.set_xlim(0, width)
        ax.set_ylim(height, 0)
        for box in svg_box:
            left, top, width, height = box["left"].inches, box["top"].inches, box["width"].inches, box["height"].inches
            picked_color = random.choice(_color_palette)
            color = (picked_color[0] / 255, picked_color[1] / 255, picked_color[2] / 255, 0.3)
            ax.add_patch(matplotlib.patches.Rectangle((left, top), width, height, color=color))
        ax.set_axis_off()
        matplotlib.pyplot.subplots_adjust(left=0, right=1, top=1, bottom=0)
        matplotlib.pyplot.margins(0)
        ax.set_xmargin(0)
        ax.set_ymargin(0)
        fig.savefig(svg_filename, bbox_inches="tight", pad_inches=0)

    svg_box = []
    _to_svg_box(slide.shapes, svg_box)
    render_svg(svg_box)

    return svg_box


def delete_shape(shape):
    """
    Delete the given shape.
    """
    parent = shape.element.getparent()
    parent.remove(shape.element)


def find_shape_with_name(shapes, name, depth=0):
    """
    Find the shape with the given name in the given shapes.
    """
    if depth == 0:
        logging.info(f"Finding shape with name: {name}")
    for shape in shapes:
        if shape.name == name:
            return shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            found = find_shape_with_name(shape.shapes, name, depth + 1)
            if found:
                return found
    return None


def find_shape_with_name_except(shapes, name, depth=0):
    """
    Find the shape with the given name in the given shapes, except the shape with the given name.
    """
    if depth == 0:
        logging.info(f"Finding shape with name: {name}")
    for shape in shapes:
        logging.debug(f"{depth * '  '} inspect shape: {shape.name}")
        if shape.name == name:
            return shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            found = find_shape_with_name(shape.shapes, name, depth + 1)
            if found:
                return found
    raise Exception(f"Shape with name {name} not found")


def _get_reltype_for_part(part):
    ctype = getattr(part, "content_type", None)
    if not ctype:
        return None
    return CONTENT_TYPE_TO_RELTYPE.get(ctype)


def _copy_and_fix_relations(original_slide, new_slide, new_el):
    """
    在 new_el 中查找所有 r:embed 和 r:id，尝试把对应的 part 从 original_slide 关联到 new_slide，
    并替换成新的 rId。返回 True/False 表示是否成功（部分失败不致命）。
    """
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    # 查找所有可能带 r:embed 或 r:id 的元素（最常见：a:blip/@r:embed, c:chart/@r:id）
    # 使用通配，遍历元素属性，找出属于 r namespace 的属性
    for el in new_el.iter():
        # el.attrib 是字典，key 可能是 '{namespace}attrname'
        for attr_name in list(el.attrib.keys()):
            if attr_name.startswith(f"{{{r_ns}}}"):
                # attr_local = attr_name.split("}", 1)[1]
                old_rId = el.attrib[attr_name]
                if not old_rId:
                    continue

                # 如果已经处理过（某些模板会多次引用同一 rId），则复用映射（我们可以缓存）
                # 先尝试从 original_slide.part 获取对应 part
                try:
                    part = original_slide.part.related_part(old_rId)
                except Exception:
                    part = None

                if part is None:
                    # 不能找到对应 part，跳过
                    continue

                # 决定新的 reltype
                reltype = _get_reltype_for_part(part)
                # fallback: 如果没映射到 reltype，使用 image 作为最后手段（保守策略）
                if reltype is None:
                    reltype = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"

                # 在 new_slide 上建立到该 part 的关系（如果之前已建立，相同 part 会创建重复关系——可用缓存优化）
                try:
                    new_rId = new_slide.part.relate_to(part, reltype)
                except Exception:
                    # 有时 relate_to 会抛异常（例如 part 类型不支持），尝试用 image reltype 作为降级
                    try:
                        new_rId = new_slide.part.relate_to(
                            part, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
                        )
                    except Exception:
                        new_rId = None

                if new_rId:
                    el.set(attr_name, new_rId)
    return True


def copy_background(source_slide, target_slide):
    source_cSld = source_slide.element.cSld
    target_cSld = target_slide.element.cSld

    # 定义命名空间
    ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}

    # 移除目标背景（bg或bgPr）
    for bg in target_cSld.findall(".//p:bg", namespaces=ns):
        target_cSld.remove(bg)
    for bgPr in target_cSld.findall(".//p:bgPr", namespaces=ns):
        target_cSld.remove(bgPr)

    # 复制源背景（必须插入到spTree之前）
    source_bg = source_cSld.find(".//p:bg", namespaces=ns)
    if source_bg is not None:
        spTree = target_cSld.find(".//p:spTree", namespaces=ns)
        target_cSld.insert(target_cSld.index(spTree), deepcopy(source_bg))


def duplicate_slide(prs, slide):
    slide_layout = slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    copy_background(slide, new_slide)

    for shape in slide.shapes:
        el = shape.element
        new_el = copy.deepcopy(el)

        # 先修复关系（图片/chart/ole 等）——使用 original slide 作为资源源
        _copy_and_fix_relations(slide, new_slide, new_el)

        # 将 new_el 插入 new_slide 的 spTree，插入到 extLst 之前（与之前一致）
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    return new_slide


# def duplicate_slide(prs, slide):
#     slide_layout = slide.slide_layout
#     new_slide = prs.slides.add_slide(slide_layout)

#     parent = new_slide.background._element.getparent()
#     # 找到当前background在父节点中的位置
#     index = parent.index(new_slide.background._element)
#     # 删除旧的，插入新的
#     parent.remove(new_slide.background._element)
#     parent.insert(index, copy.deepcopy(slide.background._element))

#     for shape in slide.shapes:
#         el = shape.element
#         new_el = copy.deepcopy(el)

#         # 处理图片 - 使用 python-pptx 内置命名空间
#         try:
#             blips = new_el.xpath(".//a:blip[@r:embed]")

#             for blip in blips:
#                 old_rId = blip.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")

#                 if old_rId:
#                     # 获取原始图片
#                     old_image_part = slide.part.related_part(old_rId)

#                     # 在新幻灯片中建立关系
#                     new_rId = new_slide.part.relate_to(
#                         old_image_part, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
#                     )

#                     # 更新 rId
#                     blip.set("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed", new_rId)

#         except (KeyError, AttributeError):
#             pass

#         new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

#     return new_slide


def delete_slide_range(prs, index_range):
    """delete slides in the given index range
    Args:
        prs: Presentation object
        index_range: range of slide indices (0-based)
    """
    for index in reversed(index_range):
        delete_slide(prs, index)


def delete_slide(prs, index):
    """delete slide at the given index
    Args:
        prs: Presentation object
        index: slide index (0-based)

    Raises:
        IndexError: when index out of range
    """
    if index < 0 or index >= len(prs.slides):
        raise IndexError(f"Slide index {index} out of range (0-{len(prs.slides) - 1})")

    xml_slides = prs.slides._sldIdLst
    xml_slides.remove(xml_slides[index])


def move_slide(prs, old_index, new_index):
    """move slide from old_index to new_index"""
    xml_slides = prs.slides._sldIdLst
    # get the element to move
    slide_element = xml_slides[old_index]
    # remove the element from old position
    xml_slides.remove(slide_element)
    # insert to new position
    xml_slides.insert(new_index, slide_element)
    return prs


def replace_picture_keep_format(slide, shape_index, new_image_path):
    shape = slide.shapes[shape_index]

    if shape.shape_type != 13:
        raise ValueError("Target shape is not a picture")

    img_id = shape._element.blip_rEmbed
    image_part = shape.part.related_part(img_id)

    with open(new_image_path, "rb") as f:
        image_part._blob = f.read()

    return shape


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser()
    parser.add_argument("--ppt", required=True)
    parser.add_argument("--method", required=True)
    args = parser.parse_args()

    from pptx import Presentation

    prs = Presentation(args.ppt)
    if args.method == "inspect":
        inspect_ppt(prs)
    elif args.method == "to_svg":
        to_svg(prs, args.ppt, "test.svg")
    else:
        raise ValueError(f"Unknown method: {args.method}")
